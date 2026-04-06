#!/usr/bin/env python3
"""
增强工具包 - 联网搜索、文件处理、API调用
"""
import json
import re
import time
from pathlib import Path
from typing import Dict, Any, List, Optional
from dataclasses import dataclass

# 尝试导入可选依赖
try:
    import requests
    HAS_REQUESTS = True
except ImportError:
    HAS_REQUESTS = False

try:
    from bs4 import BeautifulSoup
    HAS_BS4 = True
except ImportError:
    HAS_BS4 = False

try:
    from PyPDF2 import PdfReader
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


def web_search(query: str, num_results: int = 5) -> Dict[str, Any]:
    """
    联网搜索（使用 DuckDuckGo）
    
    Args:
        query: 搜索关键词
        num_results: 返回结果数量
    
    Returns:
        搜索结果字典
    """
    try:
        from duckduckgo_search import DDGS
        
        with DDGS() as ddgs:
            results = list(ddgs.text(query, max_results=num_results))
        
        return {
            "success": True,
            "query": query,
            "results": results,
            "count": len(results)
        }
    except ImportError:
        return {
            "success": False,
            "error": "缺少 duckduckgo-search 依赖，请运行: pip install duckduckgo-search"
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"搜索失败: {str(e)}"
        }


def web_fetch(url: str, max_length: int = 5000) -> Dict[str, Any]:
    """
    抓取网页内容
    
    Args:
        url: 网页 URL
        max_length: 最大返回内容长度
    
    Returns:
        网页内容字典
    """
    if not HAS_REQUESTS:
        return {
            "success": False,
            "error": "缺少 requests 依赖，请运行: pip install requests"
        }
    
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
        }
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        
        content = response.text
        
        if HAS_BS4:
            soup = BeautifulSoup(content, "html.parser")
            for script in soup(["script", "style"]):
                script.decompose()
            content = soup.get_text(separator="\n", strip=True)
        
        content = re.sub(r"\n{3,}", "\n\n", content)
        content = content.strip()
        
        if len(content) > max_length:
            content = content[:max_length] + "...\n\n[内容已截断]"
        
        return {
            "success": True,
            "url": url,
            "content": content,
            "length": len(content)
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"抓取失败: {str(e)}"
        }


def api_request(url: str, method: str = "GET", data: Optional[Dict] = None, 
                 headers: Optional[Dict] = None, timeout: int = 30) -> Dict[str, Any]:
    """
    通用 API 请求
    
    Args:
        url: API URL
        method: HTTP 方法（GET/POST/PUT/DELETE）
        data: 请求数据（JSON）
        headers: 请求头
        timeout: 超时时间（秒）
    
    Returns:
        API 响应字典
    """
    if not HAS_REQUESTS:
        return {
            "success": False,
            "error": "缺少 requests 依赖，请运行: pip install requests"
        }
    
    try:
        request_headers = headers or {}
        request_headers.setdefault("Content-Type", "application/json")
        
        if method.upper() == "GET":
            response = requests.get(url, headers=request_headers, params=data, timeout=timeout)
        elif method.upper() == "POST":
            response = requests.post(url, headers=request_headers, json=data, timeout=timeout)
        elif method.upper() == "PUT":
            response = requests.put(url, headers=request_headers, json=data, timeout=timeout)
        elif method.upper() == "DELETE":
            response = requests.delete(url, headers=request_headers, timeout=timeout)
        else:
            return {
                "success": False,
                "error": f"不支持的 HTTP 方法: {method}"
            }
        
        response.raise_for_status()
        
        try:
            result = response.json()
        except:
            result = response.text
        
        return {
            "success": True,
            "url": url,
            "method": method,
            "status_code": response.status_code,
            "result": result
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"API 请求失败: {str(e)}"
        }


def extract_pdf_text(file_path: str, max_pages: int = 50) -> Dict[str, Any]:
    """
    提取 PDF 文本内容
    
    Args:
        file_path: PDF 文件路径
        max_pages: 最大提取页数
    
    Returns:
        提取结果字典
    """
    if not HAS_PDF:
        return {
            "success": False,
            "error": "缺少 PyPDF2 依赖，请运行: pip install PyPDF2"
        }
    
    try:
        path = Path(file_path)
        if not path.exists():
            return {
                "success": False,
                "error": f"文件不存在: {file_path}"
            }
        
        reader = PdfReader(path)
        num_pages = len(reader.pages)
        pages_to_extract = min(num_pages, max_pages)
        
        text_parts = []
        for i in range(pages_to_extract):
            page = reader.pages[i]
            text_parts.append(f"--- 第 {i+1} 页 ---\n{page.extract_text()}")
        
        full_text = "\n\n".join(text_parts)
        
        return {
            "success": True,
            "file_path": file_path,
            "total_pages": num_pages,
            "extracted_pages": pages_to_extract,
            "text": full_text,
            "length": len(full_text)
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"PDF 提取失败: {str(e)}"
        }


def extract_docx_text(file_path: str) -> Dict[str, Any]:
    """
    提取 Word (.docx) 文本内容
    
    Args:
        file_path: Word 文件路径
    
    Returns:
        提取结果字典
    """
    if not HAS_DOCX:
        return {
            "success": False,
            "error": "缺少 python-docx 依赖，请运行: pip install python-docx"
        }
    
    try:
        path = Path(file_path)
        if not path.exists():
            return {
                "success": False,
                "error": f"文件不存在: {file_path}"
            }
        
        doc = Document(path)
        text_parts = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells])
                if row_text.strip():
                    text_parts.append(f"| {row_text} |")
        
        full_text = "\n\n".join(text_parts)
        
        return {
            "success": True,
            "file_path": file_path,
            "text": full_text,
            "length": len(full_text)
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"Word 提取失败: {str(e)}"
        }


def extract_xlsx_text(file_path: str, max_sheets: int = 10, max_rows: int = 100) -> Dict[str, Any]:
    """
    提取 Excel (.xlsx) 文本内容
    
    Args:
        file_path: Excel 文件路径
        max_sheets: 最大提取工作表数
        max_rows: 每个工作表最大提取行数
    
    Returns:
        提取结果字典
    """
    if not HAS_XLSX:
        return {
            "success": False,
            "error": "缺少 openpyxl 依赖，请运行: pip install openpyxl"
        }
    
    try:
        path = Path(file_path)
        if not path.exists():
            return {
                "success": False,
                "error": f"文件不存在: {file_path}"
            }
        
        workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
        sheet_names = workbook.sheetnames[:max_sheets]
        
        text_parts = []
        
        for sheet_name in sheet_names:
            sheet = workbook[sheet_name]
            text_parts.append(f"=== 工作表: {sheet_name} ===")
            
            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                if row_count >= max_rows:
                    break
                
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    text_parts.append(row_text)
                
                row_count += 1
            
            if row_count >= max_rows:
                text_parts.append(f"... (已截断，超过 {max_rows} 行)")
            
            text_parts.append("")
        
        workbook.close()
        full_text = "\n".join(text_parts)
        
        return {
            "success": True,
            "file_path": file_path,
            "sheets_extracted": sheet_names,
            "total_sheets": len(workbook.sheetnames),
            "text": full_text,
            "length": len(full_text)
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"Excel 提取失败: {str(e)}"
        }


def extract_pptx_text(file_path: str, max_slides: int = 50) -> Dict[str, Any]:
    """
    提取 PowerPoint (.pptx) 文本内容
    
    Args:
        file_path: PowerPoint 文件路径
        max_slides: 最大提取幻灯片数
    
    Returns:
        提取结果字典
    """
    if not HAS_PPTX:
        return {
            "success": False,
            "error": "缺少 python-pptx 依赖，请运行: pip install python-pptx"
        }
    
    try:
        path = Path(file_path)
        if not path.exists():
            return {
                "success": False,
                "error": f"文件不存在: {file_path}"
            }
        
        prs = Presentation(path)
        num_slides = len(prs.slides)
        slides_to_extract = min(num_slides, max_slides)
        
        text_parts = []
        
        for i in range(slides_to_extract):
            slide = prs.slides[i]
            text_parts.append(f"--- 幻灯片 {i+1} ---")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
        
        full_text = "\n\n".join(text_parts)
        
        return {
            "success": True,
            "file_path": file_path,
            "total_slides": num_slides,
            "extracted_slides": slides_to_extract,
            "text": full_text,
            "length": len(full_text)
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"PowerPoint 提取失败: {str(e)}"
        }


def extract_file_text(file_path: str) -> Dict[str, Any]:
    """
    通用文件文本提取（自动识别类型）
    
    Args:
        file_path: 文件路径
    
    Returns:
        提取结果字典
    """
    path = Path(file_path)
    if not path.exists():
        return {
            "success": False,
            "error": f"文件不存在: {file_path}"
        }
    
    ext = path.suffix.lower()
    
    if ext == ".pdf":
        return extract_pdf_text(file_path)
    elif ext == ".docx":
        return extract_docx_text(file_path)
    elif ext == ".xlsx":
        return extract_xlsx_text(file_path)
    elif ext == ".pptx":
        return extract_pptx_text(file_path)
    elif ext in [".txt", ".md", ".py", ".js", ".html", ".css", ".json", ".csv"]:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
            return {
                "success": True,
                "file_path": file_path,
                "text": text,
                "length": len(text)
            }
        except Exception as e:
            return {
                "success": False,
                "error": f"读取文本文件失败: {str(e)}"
            }
    else:
        return {
            "success": False,
            "error": f"不支持的文件类型: {ext}"
        }


__all__ = [
    "web_search",
    "web_fetch",
    "api_request",
    "extract_pdf_text",
    "extract_docx_text",
    "extract_xlsx_text",
    "extract_pptx_text",
    "extract_file_text"
]
